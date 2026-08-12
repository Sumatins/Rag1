from __future__ import annotations

import io
import os
from dataclasses import dataclass

import pandas as pd
import pymupdf
from docx import Document
from dotenv import load_dotenv
from groq import Groq
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity


load_dotenv()


# ---------------------------------------------------------
# Configuration
# ---------------------------------------------------------

EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
LLM_MODEL = "llama-3.1-8b-instant"

CHUNK_SIZE = 800
CHUNK_OVERLAP = 120
TOP_K = 4


@dataclass
class DocumentChunk:
    document_name: str
    document_type: str
    source: str
    text: str


class MultiDocumentRAG:

    def __init__(self):

        api_key = os.getenv("GROQ_API_KEY")

        if not api_key:
            raise ValueError(
                "GROQ_API_KEY not found in .env"
            )

        self.client = Groq(
            api_key=api_key
        )

        self.embedder = SentenceTransformer(
            EMBEDDING_MODEL
        )

        self.chunks: list[DocumentChunk] = []
        self.embeddings = None

    # =====================================================
    # PDF
    # =====================================================

    def extract_pdf(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        document = pymupdf.open(
            stream=file_bytes,
            filetype="pdf",
        )

        results = []

        try:

            for page_index, page in enumerate(
                document
            ):

                text = page.get_text(
                    "text"
                ).strip()

                if text:

                    results.append(
                        DocumentChunk(
                            document_name=filename,
                            document_type="PDF",
                            source=f"Page {page_index + 1}",
                            text=text,
                        )
                    )

        finally:

            document.close()

        return results

    # =====================================================
    # TXT
    # =====================================================

    def extract_txt(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        text = file_bytes.decode(
            "utf-8",
            errors="ignore",
        )

        return [
            DocumentChunk(
                document_name=filename,
                document_type="TXT",
                source="Text file",
                text=text,
            )
        ]

    # =====================================================
    # DOCX
    # =====================================================

    def extract_docx(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        document = Document(
            io.BytesIO(file_bytes)
        )

        paragraphs = []

        for paragraph in document.paragraphs:

            text = paragraph.text.strip()

            if text:

                paragraphs.append(text)

        # Extract tables too
        for table_index, table in enumerate(
            document.tables
        ):

            rows = []

            for row in table.rows:

                cells = [
                    cell.text.strip()
                    for cell in row.cells
                ]

                rows.append(
                    " | ".join(cells)
                )

            if rows:

                paragraphs.append(
                    f"Table {table_index + 1}:\n"
                    + "\n".join(rows)
                )

        text = "\n\n".join(paragraphs)

        return [
            DocumentChunk(
                document_name=filename,
                document_type="DOCX",
                source="Word document",
                text=text,
            )
        ]

    # =====================================================
    # CSV
    # =====================================================

    def extract_csv(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        dataframe = pd.read_csv(
            io.BytesIO(file_bytes)
        )

        # Convert dataframe to readable text
        text = dataframe.to_string(
            index=False
        )

        return [
            DocumentChunk(
                document_name=filename,
                document_type="CSV",
                source="CSV table",
                text=text,
            )
        ]

    # =====================================================
    # PPTX
    # =====================================================

    def extract_pptx(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        presentation = Presentation(
            io.BytesIO(file_bytes)
        )

        results = []

        for slide_index, slide in enumerate(
            presentation.slides
        ):

            slide_text = []

            for shape in slide.shapes:

                if hasattr(
                    shape,
                    "text",
                ):

                    text = shape.text.strip()

                    if text:

                        slide_text.append(
                            text
                        )

            combined_text = "\n".join(
                slide_text
            ).strip()

            if combined_text:

                results.append(
                    DocumentChunk(
                        document_name=filename,
                        document_type="PPTX",
                        source=(
                            f"Slide "
                            f"{slide_index + 1}"
                        ),
                        text=combined_text,
                    )
                )

        return results

    # =====================================================
    # Detect document type
    # =====================================================

    def extract_document(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> list[DocumentChunk]:

        extension = (
            filename
            .lower()
            .split(".")[-1]
        )

        if extension == "pdf":

            return self.extract_pdf(
                file_bytes,
                filename,
            )

        if extension == "txt":

            return self.extract_txt(
                file_bytes,
                filename,
            )

        if extension == "docx":

            return self.extract_docx(
                file_bytes,
                filename,
            )

        if extension == "csv":

            return self.extract_csv(
                file_bytes,
                filename,
            )

        if extension == "pptx":

            return self.extract_pptx(
                file_bytes,
                filename,
            )

        raise ValueError(
            f"Unsupported file type: .{extension}"
        )

    # =====================================================
    # Chunk text
    # =====================================================

    def chunk_text(
        self,
        text: str,
    ) -> list[str]:

        text = text.strip()

        if not text:
            return []

        chunks = []

        start = 0

        while start < len(text):

            end = start + CHUNK_SIZE

            chunk = text[start:end]

            chunk = chunk.strip()

            if chunk:

                chunks.append(chunk)

            if end >= len(text):
                break

            start = (
                end
                - CHUNK_OVERLAP
            )

        return chunks

    # =====================================================
    # Add document
    # =====================================================

    def add_document(
        self,
        file_bytes: bytes,
        filename: str,
    ):

        raw_chunks = self.extract_document(
            file_bytes,
            filename,
        )

        expanded_chunks = []

        for item in raw_chunks:

            pieces = self.chunk_text(
                item.text
            )

            for piece_index, piece in enumerate(
                pieces
            ):

                expanded_chunks.append(
                    DocumentChunk(
                        document_name=(
                            item.document_name
                        ),
                        document_type=(
                            item.document_type
                        ),
                        source=(
                            f"{item.source}, "
                            f"Chunk {piece_index + 1}"
                        ),
                        text=piece,
                    )
                )

        self.chunks.extend(
            expanded_chunks
        )

        return len(expanded_chunks)

    # =====================================================
    # Build embedding index
    # =====================================================

    def build_index(self):

        if not self.chunks:

            raise ValueError(
                "No document content found."
            )

        texts = [
            (
                f"{chunk.document_name}\n"
                f"{chunk.source}\n"
                f"{chunk.text}"
            )
            for chunk in self.chunks
        ]

        self.embeddings = (
            self.embedder.encode(
                texts,
                normalize_embeddings=True,
                show_progress_bar=False,
            )
        )

    # =====================================================
    # Retrieve
    # =====================================================

    def retrieve(
        self,
        question: str,
        top_k: int = TOP_K,
    ):

        if (
            self.embeddings is None
            or not self.chunks
        ):

            raise ValueError(
                "Index has not been built."
            )

        query_embedding = (
            self.embedder.encode(
                [question],
                normalize_embeddings=True,
                show_progress_bar=False,
            )
        )

        scores = cosine_similarity(
            query_embedding,
            self.embeddings,
        )[0]

        indices = sorted(
            range(len(scores)),
            key=lambda i: scores[i],
            reverse=True,
        )[:top_k]

        results = []

        for index in indices:

            results.append(
                (
                    self.chunks[index],
                    float(scores[index]),
                )
            )

        return results

    # =====================================================
    # Generate answer
    # =====================================================

    def answer(
        self,
        question: str,
        retrieved,
    ):

        context_parts = []

        for chunk, score in retrieved:

            context_parts.append(
                (
                    f"Document: "
                    f"{chunk.document_name}\n"
                    f"Type: "
                    f"{chunk.document_type}\n"
                    f"Source: "
                    f"{chunk.source}\n"
                    f"Content:\n"
                    f"{chunk.text}"
                )
            )

        context = "\n\n---\n\n".join(
            context_parts
        )

        prompt = f"""
You are a document question-answering assistant.

Answer the user's question using ONLY the
retrieved document content below.

USER QUESTION:
{question}

RETRIEVED DOCUMENT CONTENT:
{context}

Rules:

1. Do not invent information.
2. Give a clear, direct answer.
3. Use exact numbers when present.
4. For CSV data, carefully preserve column meaning.
5. For PPTX, mention the slide number.
6. For PDF, mention the page number.
7. For DOCX, mention the relevant section/table.
8. For TXT, mention the text file.
9. If the answer is not supported by the
   retrieved content, say:

"The uploaded documents do not contain
enough information to answer this question."

At the end provide:

Sources:
- document name
- document type
- source location
"""

        response = self.client.chat.completions.create(
            model=LLM_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a factual "
                        "document RAG assistant."
                    ),
                },
                {
                    "role": "user",
                    "content": prompt,
                },
            ],
            temperature=0,
            max_completion_tokens=1000,
        )

        return (
            response
            .choices[0]
            .message
            .content
            .strip()
        )